import sys
import os
import subprocess

# Auto-install python-pptx if needed
try:
    from pptx import Presentation
except ImportError:
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy
from lxml import etree

# ─── DESIGN SYSTEM ──────────────────────────────────────────────────────────────
C_NIGHT    = RGBColor(0x04, 0x34, 0x2C)   # #04342C deep bg
C_DEEP     = RGBColor(0x08, 0x50, 0x41)   # #085041 forest
C_PRIMARY  = RGBColor(0x1D, 0x9E, 0x75)   # #1D9E75 emerald
C_LIGHT    = RGBColor(0x5D, 0xCA, 0xA5)   # #5DCAA5 mint
C_PALE     = RGBColor(0x9F, 0xE1, 0xCB)   # #9FE1CB pale mint
C_BG       = RGBColor(0xE1, 0xF5, 0xEE)   # #E1F5EE light bg
C_ACCENT   = RGBColor(0xEF, 0x9F, 0x27)   # #EF9F27 safran
C_ACCENT2  = RGBColor(0xD9, 0x8A, 0x1C)   # deeper safran
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_BODY     = RGBColor(0x2C, 0x2C, 0x2A)   # #2C2C2A charcoal
C_MUTED    = RGBColor(0x5C, 0x5C, 0x57)   # muted gray
C_DANGER   = RGBColor(0xEF, 0x44, 0x44)   # red
C_GOLD     = RGBColor(0xFF, 0xD7, 0x00)   # gold

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


# ─── LOW-LEVEL XML HELPERS ──────────────────────────────────────────────────────

def rgb_hex(c: RGBColor) -> str:
    return f"{c[0]:02X}{c[1]:02X}{c[2]:02X}"


def solid_fill_xml(color: RGBColor):
    """Return lxml element for a solid fill."""
    solidFill = etree.SubElement(etree.Element("dummy"), qn("a:solidFill"))
    srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgbClr.set("val", rgb_hex(color))
    return solidFill


def grad_fill_xml(stops):
    """
    stops = list of (position_pct_int, RGBColor)
    Returns an a:gradFill element as lxml element.
    """
    gradFill = etree.Element(qn("a:gradFill"))
    gsLst = etree.SubElement(gradFill, qn("a:gsLst"))
    for pos, color in stops:
        gs = etree.SubElement(gsLst, qn("a:gs"))
        gs.set("pos", str(pos * 1000))
        srgbClr = etree.SubElement(gs, qn("a:srgbClr"))
        srgbClr.set("val", rgb_hex(color))
    lin = etree.SubElement(gradFill, qn("a:lin"))
    lin.set("ang", "5400000")  # 90 degrees (top to bottom)
    lin.set("scaled", "0")
    return gradFill


def set_shape_gradient(shape, stops):
    """Apply a gradient fill to a shape via its spPr."""
    sp = shape._element
    spPr = sp.find(qn("p:spPr"))
    if spPr is None:
        spPr = etree.SubElement(sp, qn("p:spPr"))
    # Remove existing fill
    for tag in [qn("a:solidFill"), qn("a:gradFill"), qn("a:noFill")]:
        el = spPr.find(tag)
        if el is not None:
            spPr.remove(el)
    spPr.append(grad_fill_xml(stops))
    # Remove line
    ln = spPr.find(qn("a:ln"))
    if ln is not None:
        spPr.remove(ln)
    noFillLn = etree.SubElement(etree.SubElement(etree.Element("x"), qn("a:ln")), qn("a:noFill"))
    lnEl = etree.SubElement(spPr, qn("a:ln"))
    lnEl.append(noFillLn.getparent().find(qn("a:noFill")))


def set_no_line(shape):
    shape.line.fill.background()


def set_slide_gradient(slide, stops):
    """Set slide background to gradient using XML."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = stops[0][1]
    # Override with gradient via XML
    bgPr = bg._element.find(qn("p:bgPr"))
    if bgPr is None:
        bgPr = etree.SubElement(bg._element, qn("p:bgPr"))
    for tag in [qn("a:solidFill"), qn("a:gradFill"), qn("a:noFill")]:
        el = bgPr.find(tag)
        if el is not None:
            bgPr.remove(el)
    bgPr.insert(0, grad_fill_xml(stops))


# ─── REUSABLE PRIMITIVES ────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, color, radius=None):
    """Solid filled rectangle / rounded rectangle."""
    if radius:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
        # Set rounding via adj
        sp = shape._element
        spPr = sp.find(qn("p:spPr"))
        prstGeom = spPr.find(qn("a:prstGeom"))
        if prstGeom is not None:
            avLst = prstGeom.find(qn("a:avLst"))
            if avLst is None:
                avLst = etree.SubElement(prstGeom, qn("a:avLst"))
            else:
                for gd in avLst.findall(qn("a:gd")):
                    avLst.remove(gd)
            gd = etree.SubElement(avLst, qn("a:gd"))
            gd.set("name", "adj")
            gd.set("fmla", f"val {radius}")
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    set_no_line(shape)
    return shape


def add_grad_rect(slide, l, t, w, h, stops, radius=None):
    if radius:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = stops[0][1]
    set_shape_gradient(shape, stops)
    return shape


def add_circle(slide, cx, cy, d, color):
    """Circle centered at cx, cy with diameter d."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - d/2, cy - d/2, d, d)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    set_no_line(shape)
    return shape


def tb(slide, l, t, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT,
       wrap=True, italic=False, font="Segoe UI"):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    return txBox


def tb_multi(slide, l, t, w, h, lines, size, color, bold=False,
             align=PP_ALIGN.LEFT, spacing=Pt(2), font="Segoe UI"):
    """lines = list of (text, bold_override, color_override, size_override)"""
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    first = True
    for item in lines:
        if isinstance(item, str):
            text, b, c, s = item, bold, color, size
        else:
            text = item[0]
            b    = item[1] if len(item) > 1 else bold
            c    = item[2] if len(item) > 2 else color
            s    = item[3] if len(item) > 3 else size
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = spacing
        run = p.add_run()
        run.text = text
        run.font.size = Pt(s)
        run.font.color.rgb = c
        run.font.bold = b
        run.font.name = font
    return txBox


def bullet_card(slide, l, t, w, h, title, bullets,
                bg=C_WHITE, title_color=C_DEEP, bullet_color=C_BODY,
                title_size=17, bullet_size=13.5, padding=Inches(0.28),
                card_radius=9000, accent_bar=True, accent_color=C_PRIMARY):
    """Render a glass-style card with title + bullet list."""
    add_rect(slide, l, t, w, h, bg, radius=card_radius)
    if accent_bar:
        add_rect(slide, l, t, Inches(0.055), h, accent_color, radius=0)
    bx = l + padding
    bw = w - padding - Inches(0.15)
    ty = t + padding * 0.6
    if title:
        tb(slide, bx, ty, bw, Inches(0.42), title,
           title_size, title_color, bold=True)
        cy = ty + Inches(0.47)
    else:
        cy = ty
    txBox = slide.shapes.add_textbox(bx, cy, bw, h - (cy - t) - padding * 0.6)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    for i, bul in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(5)
        dot = p.add_run()
        dot.text = "▸  "
        dot.font.color.rgb = accent_color
        dot.font.size = Pt(bullet_size - 1)
        dot.font.name = "Segoe UI"
        run = p.add_run()
        run.text = bul
        run.font.size = Pt(bullet_size)
        run.font.color.rgb = bullet_color
        run.font.name = "Segoe UI"


def divider(slide, y, color=C_PRIMARY, opacity_pct=30):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.6), y, Inches(12.133), Inches(0.012))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    set_no_line(line)


def page_label(slide, text, color=C_MUTED):
    tb(slide, Inches(12.0), Inches(7.15), Inches(1.2), Inches(0.28),
       text, 10, color, align=PP_ALIGN.RIGHT)


def slide_header(slide, title, subtitle=None, dark=False):
    tc = C_WHITE if dark else C_DEEP
    sc = C_PALE if dark else C_MUTED
    tb(slide, Inches(0.6), Inches(0.38), Inches(12.0), Inches(0.6),
       title, 28, tc, bold=True)
    if subtitle:
        tb(slide, Inches(0.6), Inches(0.98), Inches(10.0), Inches(0.35),
           subtitle, 14.5, sc, italic=True)
    # accent bar under title
    color = C_ACCENT if dark else C_PRIMARY
    add_rect(slide, Inches(0.6), Inches(0.95) if not subtitle else Inches(1.33),
             Inches(1.8), Inches(0.045), color)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — HERO / TITLE (full dark gradient)
# ═══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank)
set_slide_gradient(s1, [(0, C_NIGHT), (60, C_DEEP), (100, C_PRIMARY)])

# Diagonal decorative blocks (top-right corner)
add_grad_rect(s1, Inches(10.5), Inches(-0.3), Inches(3.8), Inches(3.8),
              [(0, C_PRIMARY), (100, C_NIGHT)], radius=6000)
add_rect(s1, Inches(11.5), Inches(0.3), Inches(1.8), Inches(1.8),
         C_ACCENT, radius=8000)

# Bottom decorative strip
add_rect(s1, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), C_DEEP)
add_rect(s1, Inches(0), Inches(7.1), Inches(3.5), Inches(0.4), C_PRIMARY)
add_rect(s1, Inches(3.5), Inches(7.1), Inches(1.0), Inches(0.4), C_ACCENT)

# Left accent bar
add_rect(s1, Inches(0), Inches(0), Inches(0.18), Inches(7.5), C_PRIMARY)

# Platform tag
add_rect(s1, Inches(0.7), Inches(1.1), Inches(3.2), Inches(0.42), C_PRIMARY, radius=21000)
tb(s1, Inches(0.85), Inches(1.15), Inches(2.9), Inches(0.32),
   "AGRITECH D'AFRIQUE DE L'OUEST", 9.5, C_WHITE, bold=True, align=PP_ALIGN.CENTER)

# Main title
tb(s1, Inches(0.7), Inches(1.72), Inches(9.0), Inches(1.6),
   "ANDD BAAY", 72, C_WHITE, bold=True, font="Segoe UI Black")

# Tagline
tb(s1, Inches(0.72), Inches(3.35), Inches(8.5), Inches(0.55),
   "L'agriculture de demain", 26, C_PALE, italic=True)

# Separator
add_rect(s1, Inches(0.72), Inches(3.98), Inches(4.5), Inches(0.055), C_ACCENT)

# Subtitle
tb(s1, Inches(0.72), Inches(4.18), Inches(8.5), Inches(0.5),
   "STRATÉGIE DE LANCEMENT, TARIFICATION & EXPÉRIENCE UTILISATEUR (UX)", 17, C_ACCENT, bold=True)

# Meta info pills
for i, (label, val) in enumerate([("Marché", "Sénégal · Mali · Sahel"),
                                    ("Modèle", "Freemium + SaaS + B2B"),
                                    ("Statut", "Prêt pour la production")]):
    x = Inches(0.72 + i * 3.3)
    add_rect(s1, x, Inches(5.1), Inches(3.1), Inches(0.55),
             RGBColor(0x0D, 0x6A, 0x52), radius=27000)
    tb(s1, x + Inches(0.12), Inches(5.14), Inches(1.0), Inches(0.4),
       label.upper(), 8.5, C_PALE, bold=True)
    tb(s1, x + Inches(1.05), Inches(5.14), Inches(1.95), Inches(0.4),
       val, 10, C_WHITE)

# Slide number
page_label(s1, "01 / 08", C_PALE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — EXECUTIVE SUMMARY (light, 3 KPI pillars)
# ═══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank)
set_slide_gradient(s2, [(0, C_BG), (100, C_WHITE)])
slide_header(s2, "Synthèse Opérationnelle",
             "Une plateforme agritech bimode et adaptative pour l'écosystème agricole sahélien.")

# 3 KPI cards across the middle
kpis = [
    (C_DEEP,    "🌍", "Marchés Cibles",   "Sénégal, Mali &\nSahel Ouest-Africain", C_PALE),
    (C_PRIMARY, "👥", "Double Interface",   "Baay Simple (ouvriers)\nBaay Pro (gestionnaires)", C_WHITE),
    (C_ACCENT,  "💳", "Modèle de Revenu",    "Gratuit · 5 000 FCFA/mois\n500 FCFA/producteur B2B", C_WHITE),
]
for i, (bg, icon, ttl, val, tc) in enumerate(kpis):
    x = Inches(0.6 + i * 4.28)
    add_rect(s2, x, Inches(1.65), Inches(3.85), Inches(2.1), bg, radius=10000)
    tb(s2, x + Inches(0.2), Inches(1.78), Inches(3.4), Inches(0.6), icon + "  " + ttl,
       17, tc, bold=True)
    tb(s2, x + Inches(0.2), Inches(2.4), Inches(3.4), Inches(0.9), val, 14.5, tc)

# Two narrative columns below
left_bullets = [
    "Statut plateforme : 43 migrations, 93/93 tests passés — déclarée prête pour la production.",
    "Moteur principal : Django 5 + WebSockets + Prédiction de rendement Scikit-Learn en direct.",
    "Système de design : Glassmorphisme émeraude, conforme WCAG AA, modes sombre & clair.",
    "Compatible PWA : Installable sur mobile, fonctionne hors ligne via les service workers.",
]
right_bullets = [
    "Le problème : Un ensemble de fonctionnalités complexes risque d'aliéner les ouvriers agricoles peu alphabétisés.",
    "La solution : Interface utilisateur adaptative selon le rôle, affichant automatiquement le mode Simple ou Pro.",
    "Stratégie GTM : Pilote → Sensibilisation physique sur les marchés (Loumas) → Expansion B2B ONG & coopératives.",
    "Monétisation : L'offre gratuite finance la croissance ; le SaaS Pro et le B2B coopératif génèrent les revenus.",
]
bullet_card(s2, Inches(0.6), Inches(4.0), Inches(5.85), Inches(3.0),
            "Où nous en sommes aujourd'hui", left_bullets,
            bg=C_WHITE, accent_color=C_PRIMARY)
bullet_card(s2, Inches(6.9), Inches(4.0), Inches(5.85), Inches(3.0),
            "L'opportunité stratégique", right_bullets,
            bg=C_WHITE, accent_color=C_ACCENT)

page_label(s2, "02 / 08")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE CORE CHALLENGE (dark left / light right split)
# ═══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank)
# Left dark panel
add_rect(s3, Inches(0), Inches(0), Inches(6.4), SLIDE_H, C_NIGHT)
add_grad_rect(s3, Inches(0), Inches(0), Inches(6.4), SLIDE_H,
              [(0, C_NIGHT), (100, C_DEEP)])
# Right light panel
add_rect(s3, Inches(6.4), Inches(0), Inches(6.933), SLIDE_H, C_BG)

# Center divider line with circle node
add_rect(s3, Inches(6.25), Inches(0), Inches(0.08), SLIDE_H, C_PRIMARY)
add_circle(s3, Inches(6.29), Inches(3.75), Inches(0.62), C_PRIMARY)
tb(s3, Inches(6.29) - Inches(0.31), Inches(3.75) - Inches(0.31) + Inches(0.06),
   Inches(0.62), Inches(0.35), "VS", 12, C_WHITE, bold=True, align=PP_ALIGN.CENTER)

# Left: field workers
tb(s3, Inches(0.45), Inches(0.45), Inches(5.7), Inches(0.6),
   "LA FRACTURE DU NUMÉRIQUE ET DE L'ALPHABÉTISATION", 13, C_PALE, bold=True)
add_rect(s3, Inches(0.45), Inches(0.95), Inches(1.6), Inches(0.045), C_ACCENT)

tb(s3, Inches(0.45), Inches(1.25), Inches(5.6), Inches(0.6),
   "Ouvriers Agricoles (Ouvriers)", 26, C_WHITE, bold=True)
tb(s3, Inches(0.45), Inches(1.88), Inches(5.6), Inches(0.4),
   "La majorité sous-servie — qualifiée sur le terrain, faible aisance numérique.", 13.5, C_PALE, italic=True)

challenges = [
    ("🗣", "Culture Orale", "La plupart parlent Wolof, Bambara ou Pulaar natively — la lecture du français est difficile."),
    ("☀", "Environnements Difficiles", "Travail en extérieur sous le soleil sur des appareils Android à petit écran d'entrée de gamme."),
    ("📵", "Faible Connectivité", "Parcelles agricoles isolées avec un signal de données mobiles faible ou nul."),
    ("😰", "Surcharge de Fonctionnalités", "Submergés par des tableaux de bord financiers denses et une navigation complexe."),
]
for i, (ic, ttl, desc) in enumerate(challenges):
    y = Inches(2.5 + i * 1.12)
    add_rect(s3, Inches(0.45), y, Inches(5.6), Inches(1.0),
             RGBColor(0x0D, 0x4A, 0x3A), radius=8000)
    tb(s3, Inches(0.6), y + Inches(0.1), Inches(0.55), Inches(0.5), ic, 20, C_ACCENT)
    tb(s3, Inches(1.2), y + Inches(0.08), Inches(4.7), Inches(0.34), ttl, 14, C_ACCENT, bold=True)
    tb(s3, Inches(1.2), y + Inches(0.42), Inches(4.7), Inches(0.5), desc, 12.5, C_PALE)

# Right: managers
tb(s3, Inches(6.65), Inches(0.45), Inches(6.1), Inches(0.6),
   "Propriétaires & Gérants", 26, C_DEEP, bold=True)
tb(s3, Inches(6.65), Inches(1.08), Inches(6.1), Inches(0.4),
   "Professionnels du numérique exigeant un contrôle approfondi des données.", 13.5, C_MUTED, italic=True)

manager_items = [
    ("📊", "Précision Financière", "Besoin de suivis de dépenses en FCFA, plafonds budgétaires et calculs de ROI en direct."),
    ("🔮", "Analyses Prédictives", "S'appuient sur Scikit-learn pour prévoir les rendements agricoles saisonniers."),
    ("💬", "Gestion d'Équipe", "Require instant WebSocket messaging, task delegation & progress tracking."),
    ("🗺", "Aperçu Spatial", "Cartes Leaflet multi-fermes, agrégation régionale des performances."),
]
for i, (ic, ttl, desc) in enumerate(manager_items):
    y = Inches(1.82 + i * 1.12)
    add_rect(s3, Inches(6.65), y, Inches(6.1), Inches(1.0), C_WHITE, radius=8000)
    add_rect(s3, Inches(6.65), y, Inches(0.055), Inches(1.0), C_PRIMARY)
    tb(s3, Inches(6.8), y + Inches(0.1), Inches(0.55), Inches(0.5), ic, 20, C_PRIMARY)
    tb(s3, Inches(7.45), y + Inches(0.08), Inches(5.2), Inches(0.34), ttl, 14, C_DEEP, bold=True)
    tb(s3, Inches(7.45), y + Inches(0.42), Inches(5.2), Inches(0.5), desc, 12.5, C_BODY)

page_label(s3, "03 / 08", C_PALE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — THE SOLUTION: DUAL ADAPTIVE UI (gradient bg)
# ═══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank)
set_slide_gradient(s4, [(0, C_NIGHT), (45, C_DEEP), (100, RGBColor(0x06, 0x55, 0x44))])

slide_header(s4, "La Solution : Double Interface Adaptative selon le Rôle",
             "Une seule plateforme. Deux expériences intelligemment adaptées et fournies automatiquement selon le rôle.", dark=True)

# Central flow arrow / label
add_rect(s4, Inches(5.8), Inches(2.05), Inches(1.73), Inches(0.52), C_PRIMARY, radius=26000)
tb(s4, Inches(5.82), Inches(2.1), Inches(1.7), Inches(0.4),
   "RÔLE DÉTECTÉ", 9, C_WHITE, bold=True, align=PP_ALIGN.CENTER)

# ── LEFT CARD: Baay Simple ──────────────────────────────────────────────────
add_grad_rect(s4, Inches(0.5), Inches(1.6), Inches(5.1), Inches(5.5),
              [(0, RGBColor(0x08, 0x70, 0x58)), (100, C_DEEP)], radius=12000)

tb(s4, Inches(0.75), Inches(1.85), Inches(4.6), Inches(0.42),
   "🌿  BAAY SIMPLE", 13.5, C_PALE, bold=True)
tb(s4, Inches(0.75), Inches(2.3), Inches(4.6), Inches(0.75),
   "L'interface Allégée", 28, C_WHITE, bold=True)
tb(s4, Inches(0.75), Inches(3.05), Inches(4.6), Inches(0.42),
   "Activée automatiquement pour les ouvriers et les producteurs peu alphabétisés.", 13, C_PALE, italic=True)

simple_features = [
    ("🎙", "Messagerie Vocale",  "Enregistrez des instructions en Wolof/Bambara. Aucune saisie requise."),
    ("🎨", "Tâches par Icônes",      "Des symboles agricoles remplacent le texte : gouttes, feuilles, outils."),
    ("📴", "PWA Hors-Ligne D'abord",      "Fonctionne sans réseau. Synchronisation dès la reconnexion 2G."),
    ("👆", "Cibles Tactiles Géantes",    "Boutons min 48×48px, grands compteurs +/–, contraste élevé."),
]
for i, (ic, ttl, desc) in enumerate(simple_features):
    y = Inches(3.6 + i * 0.84)
    tb(s4, Inches(0.75), y, Inches(0.5), Inches(0.5), ic, 15, C_ACCENT)
    tb(s4, Inches(1.35), y + Inches(0.01), Inches(4.0), Inches(0.26), ttl, 13.5, C_WHITE, bold=True)
    tb(s4, Inches(1.35), y + Inches(0.28), Inches(4.0), Inches(0.4), desc, 12, C_PALE)

# ── RIGHT CARD: Baay Pro ────────────────────────────────────────────────────
add_grad_rect(s4, Inches(7.73), Inches(1.6), Inches(5.1), Inches(5.5),
              [(0, C_ACCENT), (100, C_ACCENT2)], radius=12000)

tb(s4, Inches(7.98), Inches(1.85), Inches(4.6), Inches(0.42),
   "⚡  BAAY PRO", 13.5, RGBColor(0xFF, 0xFF, 0xFF), bold=True)
tb(s4, Inches(7.98), Inches(2.3), Inches(4.6), Inches(0.75),
   "Le Tableau de Bord", 28, C_WHITE, bold=True)
tb(s4, Inches(7.98), Inches(3.05), Inches(4.6), Inches(0.42),
   "Activée automatiquement pour les propriétaires, gérants et agronomes.", 13, RGBColor(0xFF, 0xFF, 0xDD), italic=True)

pro_features = [
    ("💰", "Cockpit Financier",   "Budget, dépenses, analyses de ROI (FCFA) & alertes de dépassement."),
    ("🌾", "Prévisions de Rendement IA",   "Prédictions Scikit-learn avec indices de confiance."),
    ("💬", "Messagerie WebSocket","Discussion en direct, réactions emoji, accusés de lecture."),
    ("🗺", "Tableau de Bord Spatial",   "Cartes Leaflet multi-fermes, vues agrégées coopératives."),
]
for i, (ic, ttl, desc) in enumerate(pro_features):
    y = Inches(3.6 + i * 0.84)
    tb(s4, Inches(7.98), y, Inches(0.5), Inches(0.5), ic, 15, C_WHITE)
    tb(s4, Inches(8.58), y + Inches(0.01), Inches(4.0), Inches(0.26), ttl, 13.5, C_WHITE, bold=True)
    tb(s4, Inches(8.58), y + Inches(0.28), Inches(4.0), Inches(0.4), desc, 12, RGBColor(0xFF, 0xFF, 0xDD))

# Centre arrow
add_rect(s4, Inches(5.85), Inches(4.0), Inches(1.63), Inches(0.055), C_PALE)
add_circle(s4, Inches(5.43), Inches(4.025), Inches(0.22), C_PALE)
add_circle(s4, Inches(7.9), Inches(4.025), Inches(0.22), C_PALE)

page_label(s4, "04 / 08", C_PALE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — BAAY SIMPLE: UX DEEP DIVE (light theme, 3 pillars)
# ═══════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank)
set_slide_gradient(s5, [(0, C_WHITE), (100, C_BG)])
slide_header(s5, "Baay Simple : Conçu pour le Terrain",
             "UX sans courbe d'apprentissage pour les contextes extérieurs, peu alphabétisés et à faible bande passante.")

# 3-column pillar layout
pillars = [
    (C_PRIMARY, "🎙", "Priorité au Vocal",
     ["Appuyer et maintenir le bouton microphone pour les messages audio",
      "Fonctionne en Wolof, Bambara et Pulaar sans saisie de texte",
      "Un haut-parleur d'aide vocale explique chaque étiquette à haute voix",
      "Messages audio sauvegardés sur Cloudinary — lisibles hors ligne",
      "Permet aux ouvriers analphabètes de participer pleinement"]),
    (C_ACCENT, "🎨", "Langage Visuel",
     ["Icônes de culture remplaçant le français : 💧 arroser, 🌿 fertiliser, 🔪 récolter",
      "Statuts avec code couleur : Vert = OK, Jaune = Attention, Rouge = Alerte",
      "Grands compteurs numériques +/– pour éviter la saisie au clavier",
      "Cibles tactiles minimales de 48×48px pour les mains gantées ou mouillées",
      "Symboles cohérents pour toutes les catégories de tâches"]),
    (C_DEEP, "📴", "PWA Hors Ligne",
     ["Le Service Worker (sw.js) met en cache toute l'interface localement",
      "Les ouvriers valident les tâches dans les champs isolés sans signal",
      "File de synchronisation intelligente téléchargeant dès la reconnexion 2G",
      "Installable comme une application native sur l'écran d'accueil Android",
      "Résolution automatique des conflits lors de la synchro multi-ouvriers"]),
]
for i, (color, ic, ttl, buls) in enumerate(pillars):
    x = Inches(0.5 + i * 4.25)
    add_grad_rect(s5, x, Inches(1.72), Inches(3.98), Inches(5.4),
                  [(0, color), (100, C_NIGHT)], radius=10000)
    # Icon circle
    add_circle(s5, x + Inches(0.75), Inches(2.38), Inches(0.9), C_WHITE)
    tb(s5, x + Inches(0.75) - Inches(0.45), Inches(2.38) - Inches(0.28),
       Inches(0.9), Inches(0.55), ic, 22, color, align=PP_ALIGN.CENTER)
    # Title
    tb(s5, x + Inches(0.22), Inches(3.05), Inches(3.5), Inches(0.5),
       ttl, 19, C_WHITE, bold=True)
    add_rect(s5, x + Inches(0.22), Inches(3.5), Inches(1.2), Inches(0.04), C_WHITE)
    # Bullets
    txb = s5.shapes.add_textbox(x + Inches(0.22), Inches(3.62), Inches(3.5), Inches(3.3))
    tf = txb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    for j, btext in enumerate(buls):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_after = Pt(5)
        dot = p.add_run()
        dot.text = "–  "
        dot.font.color.rgb = C_PALE
        dot.font.size = Pt(12.5)
        dot.font.name = "Segoe UI"
        run = p.add_run()
        run.text = btext
        run.font.size = Pt(12.5)
        run.font.color.rgb = C_WHITE
        run.font.name = "Segoe UI"

page_label(s5, "05 / 08")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — PRICING TIERS (dark bg, 3 cards)
# ═══════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(blank)
set_slide_gradient(s6, [(0, C_NIGHT), (100, C_DEEP)])
slide_header(s6, "Tarification Stratégique : Modèle à 3 Niveaux",
             "Niveaux adaptés au Sahel basés sur le Mobile Money (Wave / Orange Money) et les flux de trésorerie saisonniers.", dark=True)

tier_data = [
    {
        "bg": RGBColor(0x0D, 0x4A, 0x3A),
        "accent": C_PRIMARY,
        "badge": "GRATUIT",
        "badge_bg": C_PRIMARY,
        "name": "Baay Basique",
        "price": "0 FCFA",
        "period": "Gratuit pour toujours",
        "tag": "Petits Exploitants Individuels",
        "features": [
            "1 Profil de ferme",
            "Jusqu'à 3 membres de ferme",
            "Interface visuelle Baay Simple",
            "Alertes météo locales de base",
            "Accès en lecture à la marketplace",
        ],
        "value": "Entonnoir GTM. Stimule l'adoption\net crée le réseau de données agricoles.",
    },
    {
        "bg": C_PRIMARY,
        "accent": C_ACCENT,
        "badge": "PRO",
        "badge_bg": C_ACCENT,
        "name": "Baay Pro",
        "price": "5 000 FCFA",
        "period": "/ mois par ferme",
        "tag": "Exploitants Commerciaux",
        "features": [
            "Jusqu'à 3 fermes",
            "Membres & rôles illimités",
            "Cockpit financier complet (FCFA)",
            "IA de prédiction Scikit-learn",
            "Discussion WebSocket en direct",
        ],
        "value": "Revenus SaaS récurrents.\nAnnuel : 45 000 FCFA (1 récolte).",
    },
    {
        "bg": RGBColor(0x0D, 0x4A, 0x3A),
        "accent": C_LIGHT,
        "badge": "B2B",
        "badge_bg": C_LIGHT,
        "name": "Baay Coopérative",
        "price": "500 FCFA",
        "period": "/ producteur / mois",
        "tag": "Coopératives · ONG · États",
        "features": [
            "Panneau d'administration dédié",
            "Modélisation globale des rendements",
            "Diffusion groupée de SMS",
            "Packs de langues vocales sur mesure",
            "Accompagnement prioritaire",
        ],
        "value": "Facturation B2B / subventions.\nFinancé par les ONG et organismes.",
    },
]

for i, tier in enumerate(tier_data):
    x = Inches(0.45 + i * 4.27)
    card_h = Inches(5.2)
    # Card base
    add_rect(s6, x, Inches(1.6), Inches(4.0), card_h, tier["bg"], radius=12000)
    # Top accent strip
    add_rect(s6, x, Inches(1.6), Inches(4.0), Inches(0.09), tier["accent"], radius=0)
    # Badge
    badge_w = Inches(0.9)
    add_rect(s6, x + Inches(4.0) - badge_w - Inches(0.18), Inches(1.72),
             badge_w, Inches(0.36), tier["badge_bg"], radius=18000)
    tb(s6, x + Inches(4.0) - badge_w - Inches(0.18), Inches(1.74),
       badge_w, Inches(0.3), tier["badge"], 8.0,
       C_NIGHT if tier["badge"] in ("GRATUIT", "B2B") else C_WHITE,
       bold=True, align=PP_ALIGN.CENTER)

    # Tier name
    tb(s6, x + Inches(0.22), Inches(1.78), Inches(3.4), Inches(0.48),
       tier["name"], 20, C_WHITE, bold=True)
    # Target
    tb(s6, x + Inches(0.22), Inches(2.25), Inches(3.4), Inches(0.34),
       tier["tag"], 11.5, tier["accent"], italic=True)
    # Price
    tb(s6, x + Inches(0.22), Inches(2.68), Inches(3.5), Inches(0.62),
       tier["price"], 30, C_WHITE, bold=True)
    tb(s6, x + Inches(0.22), Inches(3.28), Inches(3.5), Inches(0.3),
       tier["period"], 11, C_PALE)
    # Divider
    add_rect(s6, x + Inches(0.22), Inches(3.65), Inches(3.5), Inches(0.04), tier["accent"])
    # Features
    txb = s6.shapes.add_textbox(x + Inches(0.22), Inches(3.76), Inches(3.5), Inches(2.0))
    tf = txb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    for j, feat in enumerate(tier["features"]):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        dot = p.add_run()
        dot.text = "✓  "
        dot.font.color.rgb = tier["accent"]
        dot.font.size = Pt(12.5)
        dot.font.name = "Segoe UI"
        dot.font.bold = True
        run = p.add_run()
        run.text = feat
        run.font.size = Pt(12.5)
        run.font.color.rgb = C_WHITE
        run.font.name = "Segoe UI"
    # Value statement
    add_rect(s6, x + Inches(0.22), Inches(5.82), Inches(3.5), Inches(0.78),
             RGBColor(0x00, 0x00, 0x00), radius=8000)
    tb(s6, x + Inches(0.32), Inches(5.9), Inches(3.3), Inches(0.62),
       tier["value"], 11, tier["accent"])

page_label(s6, "06 / 08", C_PALE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — GO-TO-MARKET LAUNCH PLAYBOOK (light theme, 3-phase timeline)
# ═══════════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(blank)
set_slide_gradient(s7, [(0, C_BG), (100, C_WHITE)])
slide_header(s7, "Go-To-Market : Plan Sahélien en 3 Phases",
             "Stratégie de confiance sur le terrain. L'agriculture rurale exige une présence physique, pas seulement des pubs.")

# Timeline horizontal rule
add_rect(s7, Inches(0.6), Inches(3.04), Inches(12.133), Inches(0.06), C_PRIMARY)

phases = [
    {
        "num": "01",
        "color": C_PRIMARY,
        "title": "Le Pilote",
        "period": "Semaines 1 – 8",
        "x": Inches(0.6),
        "bullets": [
            "Sélectionner 5 gérants de fermes à l'aise avec la tech dans les Niayes / bassin du fleuve.",
            "Déployer Baay Pro gratuitement — observer, mesurer, ajuster quotidiennement.",
            "Co-concevoir l'interface Simple : s'asseoir avec les ouvriers pour tester icônes et audio.",
            "Collecter les données de référence sur les rendements et l'avancement des tâches.",
        ],
    },
    {
        "num": "02",
        "color": C_ACCENT,
        "title": "Sensibilisation Marchés",
        "period": "Semaines 9 – 16",
        "x": Inches(4.65),
        "bullets": [
            "Installer des stands mobiles lors des jours de marchés hebdomadaires (Loumas).",
            "Démo en direct : inscrire une ferme, recevoir une alerte météo en 60 secondes.",
            "Guider les agriculteurs pour ajouter la PWA sur l'écran d'accueil Android.",
            "S'allier avec les fédérations de producteurs d'oignons, de tomates et de riz.",
        ],
    },
    {
        "num": "03",
        "color": C_DEEP,
        "title": "Expansion B2B",
        "period": "Mois 6 et plus",
        "x": Inches(8.7),
        "bullets": [
            "Présenter les tableaux de bord coopératifs agrégés à la FAO, l'USAID, la GIZ.",
            "Proposer des espaces marketplace aux sponsors d'intrants (semences/engrais).",
            "Lancer des abonnements saisonniers de diffusion de SMS d'alerte récolte.",
            "S'étendre vers un second marché (Mali, Burkina Faso, Guinée-Bissau).",
        ],
    },
]

for ph in phases:
    # Timeline dot
    add_circle(s7, ph["x"] + Inches(1.87), Inches(3.07), Inches(0.42), ph["color"])
    tb(s7, ph["x"] + Inches(1.87) - Inches(0.21), Inches(3.01), Inches(0.42), Inches(0.34),
       ph["num"], 12, C_WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Card above
    add_rect(s7, ph["x"], Inches(1.55), Inches(3.73), Inches(1.4), ph["color"], radius=8000)
    tb(s7, ph["x"] + Inches(0.2), Inches(1.68), Inches(3.3), Inches(0.5),
       ph["title"], 20, C_WHITE, bold=True)
    tb(s7, ph["x"] + Inches(0.2), Inches(2.18), Inches(3.3), Inches(0.35),
       ph["period"], 12.5, RGBColor(0xFF, 0xFF, 0xDD))

    # Card below
    bullet_card(s7, ph["x"], Inches(3.38), Inches(3.73), Inches(3.7),
                "", ph["bullets"],
                bg=C_WHITE, accent_color=ph["color"], title_size=1, bullet_size=12.5)

page_label(s7, "07 / 08")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — ENGINEERING NEXT STEPS + CLOSING (dark)
# ═══════════════════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(blank)
set_slide_gradient(s8, [(0, C_NIGHT), (55, C_DEEP), (100, RGBColor(0x07, 0x60, 0x4E))])

# Left edge strip
add_rect(s8, Inches(0), Inches(0), Inches(0.18), SLIDE_H, C_PRIMARY)

# Decorative top-right circle
add_circle(s8, Inches(11.9), Inches(1.1), Inches(3.2), RGBColor(0x06, 0x55, 0x44))
add_circle(s8, Inches(11.9), Inches(1.1), Inches(2.0), RGBColor(0x08, 0x6A, 0x58))
add_circle(s8, Inches(11.9), Inches(1.1), Inches(1.0), C_PRIMARY)
tb(s8, Inches(11.2), Inches(0.75), Inches(1.4), Inches(0.7), "🚀", 28, C_WHITE, align=PP_ALIGN.CENTER)

# Title block
tb(s8, Inches(0.6), Inches(0.4), Inches(9.5), Inches(0.55),
   "ÉTAPES SUIVANTES IMMÉDIATES", 12.5, C_PALE, bold=True)
add_rect(s8, Inches(0.6), Inches(0.95), Inches(2.0), Inches(0.05), C_ACCENT)
tb(s8, Inches(0.6), Inches(1.12), Inches(9.5), Inches(0.8),
   "Jalons Techniques & Commerciaux", 34, C_WHITE, bold=True)

steps = [
    (C_PRIMARY, "Jalon 1", "Configurer le Routage par Rôle",
     "Mettre à jour baay/permissions.py + views.py pour que les rôles Ouvriers soient automatiquement redirigés vers /dashboard/simple/ avec le thème visuel allégé."),
    (C_ACCENT,  "Jalon 2", "Développer dashboard_simple.html",
     "Modèle mobile-first et hors-ligne : grille de tâches visuelle par icônes, compteurs +/- géants, bouton micro pour discussion dialectale et cartes colorées."),
    (C_LIGHT,   "Jalon 3", "Intégrer l'Audio dans la Messagerie",
     "Ajouter le micro interactif dans la messagerie WebSocket. Sauvegarder les audios sur Cloudinary. Lecture des dialectes Wolof/Bambara."),
    (C_GOLD,    "Jalon 4", "Intégrer les Webhooks de Mobile Money",
     "Connecter les API IPN Wave & Orange Money. Ajouter le niveau d'abonnement au modèle Ferme. Gater les vues Pro derrière la vérification du niveau de facturation."),
]

col_w = Inches(5.7)
for i, (color, badge, title, desc) in enumerate(steps):
    row = i // 2
    col = i % 2
    x = Inches(0.5 + col * 6.0)
    y = Inches(2.15 + row * 2.3)
    add_rect(s8, x, y, col_w, Inches(2.1), RGBColor(0x0A, 0x3A, 0x2E), radius=10000)
    add_rect(s8, x, y, col_w, Inches(0.07), color)
    # Badge
    add_rect(s8, x + Inches(0.2), y + Inches(0.18), Inches(1.35), Inches(0.33),
             color, radius=16000)
    tb(s8, x + Inches(0.22), y + Inches(0.2), Inches(1.31), Inches(0.28),
       badge, 9.5, C_NIGHT, bold=True, align=PP_ALIGN.CENTER)
    tb(s8, x + Inches(1.7), y + Inches(0.2), Inches(3.8), Inches(0.35),
       title, 14, C_WHITE, bold=True)
    tb(s8, x + Inches(0.2), y + Inches(0.62), col_w - Inches(0.4), Inches(1.3),
       desc, 12.5, C_PALE)

# Bottom closing tagline
add_rect(s8, Inches(0), Inches(6.85), SLIDE_W, Inches(0.65), RGBColor(0x02, 0x22, 0x1B))
tb(s8, Inches(0.4), Inches(6.91), Inches(12.5), Inches(0.4),
   "\"Andd Baay\"  —  signifie en Wolof  'Cultiver Ensemble'   |   Conçu pour le Terrain. Bâti pour l'Avenir.",
   13, C_PALE, align=PP_ALIGN.CENTER)

page_label(s8, "08 / 08", C_PALE)


# ─── SAVE ───────────────────────────────────────────────────────────────────────
output_path = r"C:\Users\HP\.gemini\antigravity\brain\d1d573a1-17fa-4df9-953a-514df92573c8\Andd_Baay_Strategy_Presentation_FR.pptx"
prs.save(output_path)
project_path = r"C:\Users\HP\Andd_baay\Andd_Baay_Strategy_Presentation_FR.pptx"
prs.save(project_path)
print("OK! Presentation saved to:\n  " + output_path + "\n  " + project_path)
